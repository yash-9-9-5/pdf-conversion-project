from flask import Flask, render_template, request, send_file, jsonify, redirect, url_for, session, flash
import os
from werkzeug.utils import secure_filename
from pdf2docx import Converter
from PyPDF2 import PdfMerger
import sqlite3
from datetime import datetime
from pptx import Presentation
import openpyxl
from pdf2image import convert_from_path
import pytesseract
from PIL import Image
import io
import platform
from werkzeug.security import generate_password_hash, check_password_hash

app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = 'uploads'
app.config['MAX_CONTENT_LENGTH'] = 16 * 1024 * 1024  # 16MB max file size
app.secret_key = 'your-secret-key-here'  # Change this to a secure secret key

# Set Tesseract path for Windows
if platform.system() == 'Windows':
    pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'

# Ensure upload directory exists
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)

def init_db():
    conn = sqlite3.connect('conversions.db')
    c = conn.cursor()
    
    # Create users table
    c.execute('''CREATE TABLE IF NOT EXISTS users
                 (id INTEGER PRIMARY KEY AUTOINCREMENT,
                  username TEXT UNIQUE NOT NULL,
                  email TEXT UNIQUE NOT NULL,
                  password TEXT NOT NULL,
                  created_at DATETIME DEFAULT CURRENT_TIMESTAMP,
                  is_premium BOOLEAN DEFAULT 0)''')
    
    # Create conversions table
    c.execute('''CREATE TABLE IF NOT EXISTS conversions
                 (id INTEGER PRIMARY KEY AUTOINCREMENT,
                  user_id INTEGER,
                  filename TEXT,
                  conversion_type TEXT,
                  timestamp DATETIME,
                  status TEXT,
                  FOREIGN KEY (user_id) REFERENCES users (id))''')
    
    conn.commit()
    conn.close()

# Login required decorator
def login_required(f):
    def decorated_function(*args, **kwargs):
        if 'user_id' not in session:
            flash('Please log in to access this page', 'error')
            return redirect(url_for('login'))
        return f(*args, **kwargs)
    decorated_function.__name__ = f.__name__
    return decorated_function

@app.route('/')
@app.route('/index')
def index():
    conversion_count = 0
    is_premium = False
    if 'user_id' in session:
        conversion_count = get_daily_conversion_count(session['user_id'])
        # Get premium status
        conn = sqlite3.connect('conversions.db')
        c = conn.cursor()
        c.execute('SELECT is_premium FROM users WHERE id = ?', (session['user_id'],))
        result = c.fetchone()
        is_premium = bool(result[0]) if result else False
        conn.close()
    
    return render_template('index.html', 
                         conversion_count=conversion_count,
                         is_premium=is_premium)

@app.route('/login', methods=['GET', 'POST'])
def login():
    if request.method == 'POST':
        username = request.form['username']
        password = request.form['password']
        
        conn = sqlite3.connect('conversions.db')
        c = conn.cursor()
        c.execute('SELECT * FROM users WHERE username = ?', (username,))
        user = c.fetchone()
        conn.close()
        
        if user and check_password_hash(user[3], password):
            session['user_id'] = user[0]
            session['username'] = user[1]
            flash('Successfully logged in!', 'success')
            return redirect(url_for('index'))
        
        flash('Invalid username or password', 'error')
    return render_template('login.html')

@app.route('/register', methods=['GET', 'POST'])
def register():
    if request.method == 'POST':
        username = request.form['username']
        email = request.form['email']
        password = request.form['password']
        
        conn = sqlite3.connect('conversions.db')
        c = conn.cursor()
        
        try:
            hashed_password = generate_password_hash(password)
            c.execute('INSERT INTO users (username, email, password) VALUES (?, ?, ?)',
                     (username, email, hashed_password))
            conn.commit()
            flash('Registration successful! Please login.', 'success')
            return redirect(url_for('login'))
        except sqlite3.IntegrityError:
            flash('Username or email already exists', 'error')
        finally:
            conn.close()
            
    return render_template('register.html')

@app.route('/logout')
def logout():
    session.clear()
    flash('You have been logged out', 'success')
    return redirect(url_for('index'))

@app.route('/convert/pdf-to-ppt', methods=['POST'])
@login_required
def convert_to_ppt():
    if 'file' not in request.files:
        return jsonify({'error': 'No file uploaded'}), 400
    
    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': 'No file selected'}), 400
    
    if not file.filename.endswith('.pdf'):
        return jsonify({'error': 'File must be a PDF'}), 400
    
    filename = secure_filename(file.filename)
    filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    file.save(filepath)
    
    try:
        # Convert PDF to images
        images = convert_from_path(filepath)
        
        # Create a new PowerPoint presentation
        prs = Presentation()
        
        # Add each page as a slide
        for image in images:
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
            # Save image to temporary file
            img_path = os.path.join(app.config['UPLOAD_FOLDER'], 'temp.png')
            image.save(img_path, 'PNG')
            
            # Add image to slide
            left = 0
            top = 0
            pic = slide.shapes.add_picture(img_path, left, top)
            
            # Clean up temporary image
            os.remove(img_path)
        
        output_path = os.path.join(app.config['UPLOAD_FOLDER'], f"{os.path.splitext(filename)[0]}.pptx")
        prs.save(output_path)
        
        # Log conversion in database
        conn = sqlite3.connect('conversions.db')
        c = conn.cursor()
        c.execute('INSERT INTO conversions (user_id, filename, conversion_type, timestamp, status) VALUES (?, ?, ?, ?, ?)',
                 (session['user_id'], filename, 'pdf-to-ppt', datetime.now(), 'success'))
        conn.commit()
        conn.close()
        
        return send_file(output_path, as_attachment=True)
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/convert/pdf-to-word', methods=['POST'])
@login_required
def convert_to_word():
    if 'file' not in request.files:
        return jsonify({'error': 'No file uploaded'}), 400
    
    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': 'No file selected'}), 400
    
    if not file.filename.endswith('.pdf'):
        return jsonify({'error': 'File must be a PDF'}), 400
    
    filename = secure_filename(file.filename)
    filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    file.save(filepath)
    
    try:
        output_path = os.path.join(app.config['UPLOAD_FOLDER'], f"{os.path.splitext(filename)[0]}.docx")
        cv = Converter(filepath)
        cv.convert(output_path)
        cv.close()
        
        # Log conversion in database
        conn = sqlite3.connect('conversions.db')
        c = conn.cursor()
        c.execute('INSERT INTO conversions (user_id, filename, conversion_type, timestamp, status) VALUES (?, ?, ?, ?, ?)',
                 (session['user_id'], filename, 'pdf-to-word', datetime.now(), 'success'))
        conn.commit()
        conn.close()
        
        return send_file(output_path, as_attachment=True)
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/convert/pdf-to-excel', methods=['POST'])
@login_required
def convert_to_excel():
    if 'file' not in request.files:
        return jsonify({'error': 'No file uploaded'}), 400
    
    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': 'No file selected'}), 400
    
    if not file.filename.endswith('.pdf'):
        return jsonify({'error': 'File must be a PDF'}), 400
    
    filename = secure_filename(file.filename)
    filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    file.save(filepath)
    
    try:
        # Verify Tesseract installation
        try:
            pytesseract.get_tesseract_version()
        except Exception as e:
            return jsonify({
                'error': 'Tesseract is not properly installed. Please ensure it is installed and the path is correct.',
                'details': str(e)
            }), 500

        # Convert PDF to images
        try:
            images = convert_from_path(filepath)
        except Exception as e:
            return jsonify({
                'error': 'Failed to convert PDF to images. Please ensure Poppler is installed.',
                'details': str(e)
            }), 500
        
        # Create a new Excel workbook
        wb = openpyxl.Workbook()
        ws = wb.active
        
        # Process each page
        for i, image in enumerate(images):
            try:
                # Convert image to text using OCR
                text = pytesseract.image_to_string(image)
                
                # Split text into rows
                rows = text.split('\n')
                
                # Write to Excel
                for row_idx, row in enumerate(rows, start=1):
                    ws.cell(row=row_idx, column=i+1, value=row)
            except Exception as e:
                return jsonify({
                    'error': 'Failed to process page with OCR.',
                    'details': str(e)
                }), 500
        
        output_path = os.path.join(app.config['UPLOAD_FOLDER'], f"{os.path.splitext(filename)[0]}.xlsx")
        wb.save(output_path)
        
        # Log conversion in database
        conn = sqlite3.connect('conversions.db')
        c = conn.cursor()
        c.execute('INSERT INTO conversions (user_id, filename, conversion_type, timestamp, status) VALUES (?, ?, ?, ?, ?)',
                 (session['user_id'], filename, 'pdf-to-excel', datetime.now(), 'success'))
        conn.commit()
        conn.close()
        
        return send_file(output_path, as_attachment=True)
    except Exception as e:
        return jsonify({
            'error': 'An unexpected error occurred during conversion.',
            'details': str(e)
        }), 500

@app.route('/merge-pdfs', methods=['POST'])
@login_required
def merge_pdfs():
    if 'file1' not in request.files or 'file2' not in request.files:
        return jsonify({'error': 'Please upload both PDF files'}), 400
    
    file1 = request.files['file1']
    file2 = request.files['file2']
    
    if file1.filename == '' or file2.filename == '':
        return jsonify({'error': 'Please select both files'}), 400
    
    if not (file1.filename.endswith('.pdf') and file2.filename.endswith('.pdf')):
        return jsonify({'error': 'Both files must be PDFs'}), 400
    
    filename1 = secure_filename(file1.filename)
    filename2 = secure_filename(file2.filename)
    filepath1 = os.path.join(app.config['UPLOAD_FOLDER'], filename1)
    filepath2 = os.path.join(app.config['UPLOAD_FOLDER'], filename2)
    
    file1.save(filepath1)
    file2.save(filepath2)
    
    try:
        merger = PdfMerger()
        merger.append(filepath1)
        merger.append(filepath2)
        
        output_path = os.path.join(app.config['UPLOAD_FOLDER'], f"merged_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf")
        merger.write(output_path)
        merger.close()
        
        # Log conversion in database
        conn = sqlite3.connect('conversions.db')
        c = conn.cursor()
        c.execute('INSERT INTO conversions (user_id, filename, conversion_type, timestamp, status) VALUES (?, ?, ?, ?, ?)',
                 (session['user_id'], f"{filename1}+{filename2}", 'merge-pdfs', datetime.now(), 'success'))
        conn.commit()
        conn.close()
        
        return send_file(output_path, as_attachment=True)
    except Exception as e:
        return jsonify({'error': str(e)}), 500

# Add this function to check conversion limits
def get_daily_conversion_count(user_id):
    conn = sqlite3.connect('conversions.db')
    c = conn.cursor()
    
    # Count conversions in the last 24 hours
    c.execute('''
        SELECT COUNT(*) FROM conversions 
        WHERE user_id = ? 
        AND timestamp > datetime('now', '-24 hours')
    ''', (user_id,))
    count = c.fetchone()[0]
    conn.close()
    return count

@app.route('/convert/<conversion_type>', methods=['POST'])
@login_required
def convert_file(conversion_type):
    # Check premium status and conversion count
    conn = sqlite3.connect('conversions.db')
    c = conn.cursor()
    
    # Get user's premium status
    c.execute('SELECT is_premium FROM users WHERE id = ?', (session['user_id'],))
    is_premium = bool(c.fetchone()[0])
    
    if not is_premium:
        # Count today's conversions
        conversion_count = get_daily_conversion_count(session['user_id'])
        if conversion_count >= 5:
            conn.close()
            flash('Daily conversion limit reached! Please upgrade to premium for unlimited conversions.', 'error')
            return redirect(url_for('premium'))
    
    if 'file' not in request.files:
        return jsonify({'error': 'No file uploaded'}), 400
    
    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': 'No file selected'}), 400
    
    if not file.filename.endswith('.pdf'):
        return jsonify({'error': 'File must be a PDF'}), 400
    
    filename = secure_filename(file.filename)
    filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    file.save(filepath)
    
    try:
        if conversion_type == 'pdf-to-ppt':
            # Convert PDF to images
            images = convert_from_path(filepath)
            
            # Create a new PowerPoint presentation
            prs = Presentation()
            
            # Add each page as a slide
            for image in images:
                slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
                # Save image to temporary file
                img_path = os.path.join(app.config['UPLOAD_FOLDER'], 'temp.png')
                image.save(img_path, 'PNG')
                
                # Add image to slide
                left = 0
                top = 0
                pic = slide.shapes.add_picture(img_path, left, top)
                
                # Clean up temporary image
                os.remove(img_path)
            
            output_path = os.path.join(app.config['UPLOAD_FOLDER'], f"{os.path.splitext(filename)[0]}.pptx")
            prs.save(output_path)
        elif conversion_type == 'pdf-to-word':
            output_path = os.path.join(app.config['UPLOAD_FOLDER'], f"{os.path.splitext(filename)[0]}.docx")
            cv = Converter(filepath)
            cv.convert(output_path)
            cv.close()
        elif conversion_type == 'pdf-to-excel':
            # Verify Tesseract installation
            try:
                pytesseract.get_tesseract_version()
            except Exception as e:
                return jsonify({
                    'error': 'Tesseract is not properly installed. Please ensure it is installed and the path is correct.',
                    'details': str(e)
                }), 500

            # Convert PDF to images
            try:
                images = convert_from_path(filepath)
            except Exception as e:
                return jsonify({
                    'error': 'Failed to convert PDF to images. Please ensure Poppler is installed.',
                    'details': str(e)
                }), 500
            
            # Create a new Excel workbook
            wb = openpyxl.Workbook()
            ws = wb.active
            
            # Process each page
            for i, image in enumerate(images):
                try:
                    # Convert image to text using OCR
                    text = pytesseract.image_to_string(image)
                    
                    # Split text into rows
                    rows = text.split('\n')
                    
                    # Write to Excel
                    for row_idx, row in enumerate(rows, start=1):
                        ws.cell(row=row_idx, column=i+1, value=row)
                except Exception as e:
                    return jsonify({
                        'error': 'Failed to process page with OCR.',
                        'details': str(e)
                    }), 500
            
            output_path = os.path.join(app.config['UPLOAD_FOLDER'], f"{os.path.splitext(filename)[0]}.xlsx")
            wb.save(output_path)
        else:
            return jsonify({'error': 'Invalid conversion type'}), 400
        
        # Log conversion in database
        conn = sqlite3.connect('conversions.db')
        c = conn.cursor()
        c.execute('INSERT INTO conversions (user_id, filename, conversion_type, timestamp, status) VALUES (?, ?, ?, ?, ?)',
                 (session['user_id'], filename, conversion_type, datetime.now(), 'success'))
        conn.commit()
        conn.close()
        
        return send_file(output_path, as_attachment=True)
    except Exception as e:
        return jsonify({
            'error': 'An unexpected error occurred during conversion.',
            'details': str(e)
        }), 500

# Add premium page route
@app.route('/premium')
@login_required
def premium():
    return render_template('premium.html')

# Add user history route
@app.route('/history')
@login_required
def history():
    conn = sqlite3.connect('conversions.db')
    c = conn.cursor()
    c.execute('''SELECT filename, conversion_type, timestamp, status 
                 FROM conversions WHERE user_id = ? 
                 ORDER BY timestamp DESC''', (session['user_id'],))
    conversions = c.fetchall()
    conn.close()
    return render_template('history.html', conversions=conversions)

if __name__ == '__main__':
    init_db()
    app.run(host='127.0.0.1', port=5002, debug=True) 